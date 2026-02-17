"""
Empire v7.3 - Document Generator Service Tests
Phase 2: Document Generation Pipeline

Tests for:
- DOCX generation (headings, tables, code blocks, lists)
- XLSX generation (styled headers, data rows, auto-width)
- PPTX generation (title slide, content slides, table slides)
- Markdown generation
- Preview generation
- Upload to B2 storage
"""

import pytest
from unittest.mock import AsyncMock, MagicMock, patch
from io import BytesIO

from app.services.document_generator_service import (
    DocumentGeneratorService,
    DocumentFormat,
    GeneratedDocument,
    MIME_TYPES,
)
from app.services.output_architect_service import ContentBlock


# ============================================================================
# Fixtures
# ============================================================================

@pytest.fixture
def service():
    return DocumentGeneratorService()


@pytest.fixture
def sample_blocks():
    """A representative set of content blocks."""
    return [
        ContentBlock(type="heading", content="Executive Summary", metadata={"level": 2}),
        ContentBlock(type="paragraph", content="Revenue grew 15% in Q4 2025."),
        ContentBlock(
            type="table",
            content="",
            metadata={
                "headers": ["Metric", "Q3", "Q4", "Change"],
                "rows": [
                    ["Revenue", "$10M", "$11.5M", "+15%"],
                    ["EBITDA", "$3M", "$3.8M", "+27%"],
                ],
            },
        ),
        ContentBlock(type="heading", content="Details", metadata={"level": 3}),
        ContentBlock(type="list", content="Item one\nItem two\nItem three", metadata={"ordered": False}),
        ContentBlock(type="code", content="SELECT * FROM revenue WHERE quarter = 'Q4';", metadata={"language": "sql"}),
    ]


@pytest.fixture
def minimal_blocks():
    """Minimal content for simple tests."""
    return [
        ContentBlock(type="heading", content="Test", metadata={"level": 2}),
        ContentBlock(type="paragraph", content="Hello world."),
    ]


# ============================================================================
# Markdown Generation
# ============================================================================

class TestMarkdownGeneration:
    @pytest.mark.asyncio
    async def test_generate_markdown(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.MARKDOWN, "Test Report")

        assert doc.format == DocumentFormat.MARKDOWN
        assert doc.mime_type == "text/markdown"
        assert doc.size_bytes > 0
        assert doc.filename.endswith(".md")

        content = doc.content_bytes.decode("utf-8")
        assert "# Test Report" in content
        assert "Executive Summary" in content
        assert "Revenue grew 15%" in content
        assert "| Metric |" in content
        assert "- Item one" in content
        assert "```sql" in content

    @pytest.mark.asyncio
    async def test_markdown_with_summary(self, service, minimal_blocks):
        doc = await service.generate(
            minimal_blocks, DocumentFormat.MARKDOWN, "Title", summary="A brief summary"
        )
        content = doc.content_bytes.decode("utf-8")
        assert "A brief summary" in content


# ============================================================================
# DOCX Generation
# ============================================================================

class TestDocxGeneration:
    @pytest.mark.asyncio
    async def test_generate_docx(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.DOCX, "Financial Report")

        assert doc.format == DocumentFormat.DOCX
        assert doc.mime_type == MIME_TYPES[DocumentFormat.DOCX]
        assert doc.size_bytes > 0
        assert doc.filename.endswith(".docx")

        # Verify it's a valid ZIP (DOCX is a ZIP container)
        import zipfile
        buf = BytesIO(doc.content_bytes)
        assert zipfile.is_zipfile(buf)

    @pytest.mark.asyncio
    async def test_docx_contains_content(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.DOCX, "Test")

        # Parse the DOCX to verify content
        from docx import Document
        docx = Document(BytesIO(doc.content_bytes))

        # Check title
        texts = [p.text for p in docx.paragraphs]
        assert any("Test" in t for t in texts)

        # Check table exists
        assert len(docx.tables) >= 1
        # Check table has correct headers
        header_cells = [cell.text for cell in docx.tables[0].rows[0].cells]
        assert "Metric" in header_cells

    @pytest.mark.asyncio
    async def test_docx_empty_blocks(self, service):
        doc = await service.generate([], DocumentFormat.DOCX, "Empty Doc")
        assert doc.size_bytes > 0  # Still produces a valid DOCX


# ============================================================================
# XLSX Generation
# ============================================================================

class TestXlsxGeneration:
    @pytest.mark.asyncio
    async def test_generate_xlsx(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.XLSX, "Revenue Data")

        assert doc.format == DocumentFormat.XLSX
        assert doc.mime_type == MIME_TYPES[DocumentFormat.XLSX]
        assert doc.filename.endswith(".xlsx")

        # Verify it's a valid ZIP (XLSX is a ZIP container)
        import zipfile
        buf = BytesIO(doc.content_bytes)
        assert zipfile.is_zipfile(buf)

    @pytest.mark.asyncio
    async def test_xlsx_contains_data(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.XLSX, "Data")

        from openpyxl import load_workbook
        wb = load_workbook(BytesIO(doc.content_bytes))
        ws = wb.active

        # Check that table data is in the sheet
        all_values = []
        for row in ws.iter_rows(values_only=True):
            all_values.extend([str(v) for v in row if v])

        assert any("Revenue" in v for v in all_values)
        assert any("$10M" in v for v in all_values)


# ============================================================================
# PPTX Generation
# ============================================================================

class TestPptxGeneration:
    @pytest.mark.asyncio
    async def test_generate_pptx(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.PPTX, "Board Presentation")

        assert doc.format == DocumentFormat.PPTX
        assert doc.mime_type == MIME_TYPES[DocumentFormat.PPTX]
        assert doc.filename.endswith(".pptx")

        # Verify it's a valid ZIP (PPTX is a ZIP container)
        import zipfile
        buf = BytesIO(doc.content_bytes)
        assert zipfile.is_zipfile(buf)

    @pytest.mark.asyncio
    async def test_pptx_has_title_slide(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.PPTX, "My Presentation")

        from pptx import Presentation
        prs = Presentation(BytesIO(doc.content_bytes))

        # First slide should be the title slide
        assert len(prs.slides) >= 1
        title_slide = prs.slides[0]
        assert title_slide.shapes.title.text == "My Presentation"

    @pytest.mark.asyncio
    async def test_pptx_heading_creates_new_slide(self, service, sample_blocks):
        doc = await service.generate(sample_blocks, DocumentFormat.PPTX, "Test")

        from pptx import Presentation
        prs = Presentation(BytesIO(doc.content_bytes))

        # Title slide + at least 2 heading slides
        assert len(prs.slides) >= 3


# ============================================================================
# Preview Generation
# ============================================================================

class TestPreviewGeneration:
    def test_preview_includes_heading(self, service, sample_blocks):
        preview = service._generate_preview(sample_blocks)
        assert "Executive Summary" in preview

    def test_preview_truncated(self, service):
        long_blocks = [
            ContentBlock(type="paragraph", content="A" * 1000)
        ]
        preview = service._generate_preview(long_blocks, max_length=100)
        assert len(preview) <= 103  # 100 + "..."

    def test_preview_includes_table_summary(self, service, sample_blocks):
        preview = service._generate_preview(sample_blocks)
        assert "Table:" in preview or "Metric" in preview

    def test_preview_empty_blocks(self, service):
        preview = service._generate_preview([])
        assert preview == ""


# ============================================================================
# Filename Sanitization
# ============================================================================

class TestFilenameSanitization:
    @pytest.mark.asyncio
    async def test_special_chars_removed(self, service, minimal_blocks):
        doc = await service.generate(minimal_blocks, DocumentFormat.MARKDOWN, "My Report! @#$%")
        assert "@" not in doc.filename
        assert "#" not in doc.filename

    @pytest.mark.asyncio
    async def test_spaces_replaced(self, service, minimal_blocks):
        doc = await service.generate(minimal_blocks, DocumentFormat.MARKDOWN, "My Big Report")
        assert " " not in doc.filename
        assert "_" in doc.filename

    @pytest.mark.asyncio
    async def test_long_title_truncated(self, service, minimal_blocks):
        doc = await service.generate(minimal_blocks, DocumentFormat.MARKDOWN, "A" * 200)
        # Filename (minus extension) should be <= 80 chars
        name_part = doc.filename.rsplit(".", 1)[0]
        assert len(name_part) <= 80


# ============================================================================
# Upload to Storage
# ============================================================================

class TestUploadToStorage:
    @pytest.mark.asyncio
    async def test_upload_sets_storage_url(self, service, minimal_blocks):
        doc = await service.generate(minimal_blocks, DocumentFormat.MARKDOWN, "Test")

        with patch("app.services.b2_storage.get_b2_service") as mock_b2_factory:
            mock_b2 = MagicMock()
            mock_b2.upload_file = AsyncMock(return_value={
                "url": "https://b2.example.com/artifacts/test.md",
                "file_name": "user-1/session-1/test.md",
            })
            mock_b2_factory.return_value = mock_b2

            updated = await service.upload_to_storage(doc, "user-1", "session-1")

            assert updated.storage_url == "https://b2.example.com/artifacts/test.md"
            assert updated.storage_path == "user-1/session-1/test.md"


# ============================================================================
# Unsupported Format
# ============================================================================

class TestUnsupportedFormat:
    @pytest.mark.asyncio
    async def test_invalid_format_raises(self, service, minimal_blocks):
        with pytest.raises(ValueError, match="Unsupported format"):
            await service.generate(minimal_blocks, "csv", "Test")  # type: ignore


# ============================================================================
# MIME Types
# ============================================================================

class TestMimeTypes:
    def test_all_formats_have_mime(self):
        for fmt in DocumentFormat:
            assert fmt in MIME_TYPES
            assert len(MIME_TYPES[fmt]) > 0

    def test_docx_mime(self):
        assert "wordprocessing" in MIME_TYPES[DocumentFormat.DOCX]

    def test_xlsx_mime(self):
        assert "spreadsheet" in MIME_TYPES[DocumentFormat.XLSX]

    def test_pptx_mime(self):
        assert "presentation" in MIME_TYPES[DocumentFormat.PPTX]
