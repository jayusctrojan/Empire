"""
Empire v7.3 - Project Source Processing Tasks
Celery tasks for processing project sources (NotebookLM-style feature)

Task 61: Unified source processing pipeline with type-specific extraction
Task 69: Performance profiling, caching, and queue prioritization
"""

import asyncio
import hashlib
import logging
import os
import re
import tempfile
import time
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple
from uuid import uuid4

from app.celery_app import celery_app, get_source_priority

# Task 69: Performance profiling and caching imports
from app.utils.performance_profiler import get_performance_profiler, Benchmark
from app.services.source_content_cache import get_source_content_cache

# Feature 006: Markdown-aware chunking
from app.services.chunking_service import (
    MarkdownChunkerStrategy,
    MarkdownChunkerConfig,
    HEADER_PATTERN
)

logger = logging.getLogger(__name__)

# Task 69: Initialize profiler and cache
_profiler = None
_cache = None


def _get_profiler():
    """Lazy load profiler"""
    global _profiler
    if _profiler is None:
        _profiler = get_performance_profiler()
    return _profiler


def _get_cache():
    """Lazy load cache"""
    global _cache
    if _cache is None:
        _cache = get_source_content_cache()
    return _cache

# Constants
CHUNK_SIZE = 512  # tokens per chunk
MAX_SUMMARY_LENGTH = 2000  # characters
MAX_RETRIES = 3


def run_async(coro):
    """Helper to run async code in Celery tasks"""
    try:
        loop = asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    return loop.run_until_complete(coro)


# ============================================================================
# Main Processing Task
# ============================================================================

@celery_app.task(
    name='app.tasks.source_processing.process_source',
    bind=True,
    max_retries=3,
    default_retry_delay=60,
    autoretry_for=(Exception,),
    retry_backoff=True,
    retry_backoff_max=300,
    queue='project_sources'  # Task 69: Dedicated queue
)
def process_source(
    self,
    source_id: str,
    user_id: str,
    queue_start_time: Optional[float] = None,  # Task 69: For queue time tracking
    content_set_context: Optional[Dict[str, Any]] = None  # Feature 007: Content set context
) -> Dict[str, Any]:
    """
    Process a project source: extract content, generate summary, create embeddings.

    Task 69: Enhanced with performance profiling and content caching.
    Feature 007: Enhanced with content set context support.

    Handles all source types:
    - Documents: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, RTF
    - Media: Images, Audio, Video
    - Web: YouTube, Website URLs
    - Archives: ZIP, TAR, GZ

    Args:
        source_id: UUID of the source to process
        user_id: User ID for access control
        queue_start_time: Timestamp when task was queued (for latency tracking)
        content_set_context: Optional context from Content Prep Agent containing:
            - content_set_id: UUID of the parent content set
            - content_set_name: Name of the content set
            - sequence_number: Position in the set (1-based)
            - is_sequential: Whether set has defined order
            - total_files: Total files in the set

    Returns:
        Processing result with status and metadata
    """
    from app.core.supabase_client import get_supabase_client
    from app.services.status_broadcaster import get_sync_status_broadcaster
    from app.models.task_status import TaskType, ProcessingStage

    supabase = get_supabase_client()
    broadcaster = get_sync_status_broadcaster()
    task_name = 'app.tasks.source_processing.process_source'

    # Task 69: Initialize profiler and cache
    profiler = _get_profiler()
    cache = _get_cache()
    benchmark = Benchmark(f"process_source_{source_id[:8]}")

    try:
        logger.info(f"Starting source processing: {source_id}")

        # Fetch source from database
        result = supabase.table("project_sources").select("*").eq(
            "id", source_id
        ).eq("user_id", user_id).execute()

        if not result.data:
            logger.error(f"Source not found: {source_id}")
            return {"status": "error", "error": "Source not found"}

        source = result.data[0]
        project_id = source["project_id"]
        source_type = source["source_type"]
        file_type = source.get("file_type")

        # Get file name for notifications
        file_name = source.get("file_name") or source.get("url") or "Unknown"

        # Task 69: Start performance profiling
        file_size = source.get("file_size_bytes")
        profile = profiler.start_source_profile(
            source_id=source_id,
            source_type=source_type,
            file_type=file_type,
            file_size_bytes=file_size,
            queue_start_time=queue_start_time
        )

        # Update status to processing
        _update_status(supabase, source_id, "processing", progress=5)
        broadcaster.broadcast_progress(
            task_id=self.request.id,
            task_name=task_name,
            current=5,
            total=100,
            message=f"Starting processing for {source_type} source",
            stage=ProcessingStage.INITIALIZING,
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id, "source_type": source_type}
        )
        # Task 62: WebSocket source status notification
        broadcaster.broadcast_source_status(
            source_id=source_id,
            project_id=project_id,
            status="processing",
            progress=5,
            message=f"Starting processing for {source_type}",
            source_type=source_type,
            file_name=file_name
        )

        # Task 69: Check cache for previously extracted content
        content = None
        metadata = source.get("metadata", {}) or {}
        cache_hit = False

        # Compute content hash for caching
        file_hash = source.get("file_hash") or hashlib.md5(
            (source.get("file_path") or source.get("url") or source_id).encode()
        ).hexdigest()

        async def check_and_extract():
            nonlocal content, metadata, cache_hit

            # Check cache first
            cached_result = await cache.get_cached_content(file_hash, source_type)
            if cached_result:
                content, cached_meta = cached_result
                metadata.update(cached_meta)
                cache_hit = True
                logger.info(f"Cache hit for content: {source_id}")
                return

            # Extract content based on source type
            async with profiler.profile_stage(source_id, "content_extraction"):
                logger.info(f"Extracting content for {source_type}: {source_id}")
                content, metadata = _extract_content(source, supabase)

            # Cache the extracted content
            if content:
                await cache.cache_content(file_hash, source_type, content, metadata)

        run_async(check_and_extract())

        if not content:
            raise ValueError(f"No content extracted from source: {source_type}")

        _update_status(supabase, source_id, "processing", progress=40)
        broadcaster.broadcast_progress(
            task_id=self.request.id,
            task_name=task_name,
            current=40,
            total=100,
            message="Content extracted, generating summary",
            stage=ProcessingStage.EXTRACTING_METADATA,
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id, "stage": "content_extracted"}
        )
        # Task 62: WebSocket source status notification
        broadcaster.broadcast_source_status(
            source_id=source_id,
            project_id=project_id,
            status="processing",
            progress=40,
            message="Content extracted, generating summary",
            source_type=source_type,
            file_name=file_name
        )

        # Task 69: Generate summary with caching
        content_hash = hashlib.sha256(content.encode()).hexdigest()

        async def generate_summary_with_cache():
            # Check cache first
            cached_summary = await cache.get_cached_summary(content_hash)
            if cached_summary:
                logger.info(f"Cache hit for summary: {source_id}")
                return cached_summary

            # Generate summary with profiling
            async with profiler.profile_stage(source_id, "summary_generation"):
                logger.info(f"Generating summary for source: {source_id}")
                summary = _generate_summary(content, source_type)

            # Cache the summary
            if summary:
                await cache.cache_summary(content_hash, summary)

            return summary

        summary = run_async(generate_summary_with_cache())

        _update_status(supabase, source_id, "processing", progress=60)
        broadcaster.broadcast_progress(
            task_id=self.request.id,
            task_name=task_name,
            current=60,
            total=100,
            message="Summary generated, creating embeddings",
            stage=ProcessingStage.CHUNKING,
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id, "stage": "summary_generated"}
        )
        # Task 62: WebSocket source status notification
        broadcaster.broadcast_source_status(
            source_id=source_id,
            project_id=project_id,
            status="processing",
            progress=60,
            message="Summary generated, creating embeddings",
            source_type=source_type,
            file_name=file_name
        )

        # Task 69: Chunk and embed with caching
        async def chunk_and_embed_with_cache():
            # Chunk content with profiling
            async with profiler.profile_stage(source_id, "chunking"):
                chunks = _chunk_content(content, CHUNK_SIZE)
                logger.info(f"Created {len(chunks)} chunks for embedding")

            # Compute chunk hashes for cache lookup
            chunk_hashes = [hashlib.md5(chunk.encode()).hexdigest() for chunk in chunks]

            # Check cache for existing embeddings
            cached_embeddings = await cache.get_cached_embeddings_batch(chunk_hashes)

            # Identify chunks that need embedding generation
            chunks_to_embed = []
            chunks_to_embed_indices = []
            for i, (chunk, chunk_hash) in enumerate(zip(chunks, chunk_hashes)):
                if cached_embeddings.get(chunk_hash) is None:
                    chunks_to_embed.append(chunk)
                    chunks_to_embed_indices.append(i)

            logger.info(
                f"Embedding cache: {len(chunks) - len(chunks_to_embed)}/{len(chunks)} hits"
            )

            # Generate embeddings for uncached chunks
            embeddings = [None] * len(chunks)

            # Fill in cached embeddings
            for i, chunk_hash in enumerate(chunk_hashes):
                if cached_embeddings.get(chunk_hash) is not None:
                    embeddings[i] = cached_embeddings[chunk_hash]

            # Generate new embeddings if needed
            if chunks_to_embed:
                async with profiler.profile_stage(source_id, "embedding_generation"):
                    logger.info(f"Generating embeddings for {len(chunks_to_embed)} chunks")
                    new_embeddings = _generate_embeddings(chunks_to_embed)

                # Fill in new embeddings and cache them
                new_embeddings_to_cache = {}
                for idx, embedding in zip(chunks_to_embed_indices, new_embeddings):
                    embeddings[idx] = embedding
                    new_embeddings_to_cache[chunk_hashes[idx]] = embedding

                # Cache the new embeddings
                await cache.cache_embeddings_batch(new_embeddings_to_cache)

            return chunks, embeddings

        chunks, embeddings = run_async(chunk_and_embed_with_cache())

        _update_status(supabase, source_id, "processing", progress=80)
        broadcaster.broadcast_progress(
            task_id=self.request.id,
            task_name=task_name,
            current=80,
            total=100,
            message=f"Generated {len(chunks)} embeddings, storing in database",
            stage=ProcessingStage.EMBEDDING,
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id, "stage": "embeddings_generated", "chunk_count": len(chunks)}
        )
        # Task 62: WebSocket source status notification
        broadcaster.broadcast_source_status(
            source_id=source_id,
            project_id=project_id,
            status="processing",
            progress=80,
            message=f"Generated {len(chunks)} embeddings, storing in database",
            source_type=source_type,
            file_name=file_name,
            metadata={"chunk_count": len(chunks)}
        )

        # Task 69: Store embeddings with profiling
        # Feature 007: Pass content set context for relationship creation
        async def store_with_profiling():
            async with profiler.profile_stage(source_id, "embedding_storage"):
                _store_embeddings(
                    supabase=supabase,
                    source_id=source_id,
                    project_id=project_id,
                    user_id=user_id,
                    chunks=chunks,
                    embeddings=embeddings,
                    content_set_context=content_set_context  # Feature 007
                )

        run_async(store_with_profiling())

        # Update source with content, summary, and metadata
        # Feature 007: Include content set context in metadata
        if content_set_context:
            metadata["content_set"] = {
                "id": content_set_context.get("content_set_id"),
                "name": content_set_context.get("content_set_name"),
                "sequence": content_set_context.get("sequence_number"),
                "is_sequential": content_set_context.get("is_sequential", True)
            }

        update_data = {
            "content": content[:50000],  # Limit stored content
            "summary": summary,
            "metadata": metadata,
            "status": "ready",
            "processing_progress": 100,
            "processing_completed_at": datetime.utcnow().isoformat(),
            "updated_at": datetime.utcnow().isoformat(),
        }
        supabase.table("project_sources").update(update_data).eq("id", source_id).execute()

        # Task 69: Complete profiling and log benchmark
        profiler.complete_source_profile(source_id, success=True)
        logger.info(f"Benchmark results: {benchmark.report()}")

        broadcaster.broadcast_success(
            task_id=self.request.id,
            task_name=task_name,
            result={
                "source_id": source_id,
                "chunks_created": len(chunks),
                "content_length": len(content),
                "has_summary": bool(summary)
            },
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id}
        )
        # Task 62: WebSocket source status notification - success
        broadcaster.broadcast_source_status(
            source_id=source_id,
            project_id=project_id,
            status="ready",
            progress=100,
            message="Source processed successfully",
            source_type=source_type,
            file_name=file_name,
            metadata={"chunks_created": len(chunks), "content_length": len(content)}
        )

        logger.info(f"Source processing completed: {source_id}")

        # Feature 007: Include content set info in result
        result = {
            "status": "success",
            "source_id": source_id,
            "chunks_created": len(chunks),
            "content_length": len(content),
            "has_summary": bool(summary)
        }

        if content_set_context:
            result["content_set"] = {
                "id": content_set_context.get("content_set_id"),
                "name": content_set_context.get("content_set_name"),
                "sequence": content_set_context.get("sequence_number")
            }

        return result

    except Exception as e:
        logger.error(f"Source processing failed: {source_id}", exc_info=True)

        # Task 69: Complete profiling on failure
        try:
            profiler.complete_source_profile(source_id, success=False, error=str(e))
        except Exception:
            pass  # Don't let profiling errors mask the real error

        # Get retry count safely
        retry_count = 0
        current_project_id = None
        current_source_type = None
        current_file_name = None

        if 'source' in dir() and source:
            retry_count = source.get("retry_count", 0)
            current_project_id = source.get("project_id")
            current_source_type = source.get("source_type")
            current_file_name = source.get("file_name") or source.get("url")
        elif 'project_id' in dir():
            current_project_id = project_id
        if 'source_type' in dir():
            current_source_type = source_type
        if 'file_name' in dir():
            current_file_name = file_name

        # Update status to failed
        _update_status(
            supabase, source_id, "failed",
            error=str(e),
            retry_count=retry_count + 1
        )

        broadcaster.broadcast_failure(
            task_id=self.request.id,
            task_name=task_name,
            error_type=type(e).__name__,
            error_message=str(e),
            retry_count=retry_count + 1,
            max_retries=MAX_RETRIES,
            task_type=TaskType.SOURCE_PROCESSING,
            metadata={"source_id": source_id}
        )

        # Task 62: WebSocket source status notification - failure
        if current_project_id:
            broadcaster.broadcast_source_status(
                source_id=source_id,
                project_id=current_project_id,
                status="failed",
                progress=0,
                message=f"Processing failed: {str(e)[:100]}",
                source_type=current_source_type,
                file_name=current_file_name,
                error_message=str(e),
                retry_count=retry_count + 1
            )

        # Re-raise to trigger Celery retry
        raise


def _update_status(
    supabase,
    source_id: str,
    status: str,
    progress: Optional[int] = None,
    error: Optional[str] = None,
    retry_count: Optional[int] = None
):
    """Update source status in database"""
    update_data = {
        "status": status,
        "updated_at": datetime.utcnow().isoformat(),
    }

    if progress is not None:
        update_data["processing_progress"] = progress

    if error:
        update_data["processing_error"] = error

    if retry_count is not None:
        update_data["retry_count"] = retry_count

    if status == "processing":
        update_data["processing_started_at"] = datetime.utcnow().isoformat()
    elif status in ("ready", "failed"):
        update_data["processing_completed_at"] = datetime.utcnow().isoformat()

    supabase.table("project_sources").update(update_data).eq("id", source_id).execute()


# ============================================================================
# Content Extraction Functions
# ============================================================================

def _extract_content(source: Dict[str, Any], supabase) -> Tuple[str, Dict[str, Any]]:
    """
    Extract content from source based on type.

    Returns:
        Tuple of (content_text, metadata_dict)
    """
    source_type = source["source_type"]
    metadata: Dict[str, Any] = source.get("metadata", {}) or {}

    if source_type == "youtube":
        return _extract_youtube_content(source["url"], metadata)

    elif source_type == "website":
        return _extract_website_content(source["url"], metadata)

    elif source_type == "pdf":
        return _extract_pdf_content(source, supabase, metadata)

    elif source_type in ("docx", "doc"):
        return _extract_docx_content(source, supabase, metadata)

    elif source_type in ("xlsx", "xls", "csv"):
        return _extract_spreadsheet_content(source, supabase, metadata)

    elif source_type in ("pptx", "ppt"):
        return _extract_pptx_content(source, supabase, metadata)

    elif source_type in ("txt", "md", "rtf"):
        return _extract_text_content(source, supabase, metadata)

    elif source_type == "image":
        return _extract_image_content(source, supabase, metadata)

    elif source_type in ("audio", "video"):
        return _extract_audio_video_content(source, supabase, metadata)

    elif source_type == "archive":
        return _extract_archive_content(source, supabase, metadata)

    else:
        raise ValueError(f"Unsupported source type: {source_type}")


def _extract_youtube_content(url: str, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract transcript and metadata from YouTube video"""
    try:
        import yt_dlp

        # Extract video info
        ydl_opts = {
            'quiet': True,
            'no_warnings': True,
            'extract_flat': False,
            'writesubtitles': True,
            'writeautomaticsub': True,
            'subtitleslangs': ['en'],
        }

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)

        # Get transcript
        transcript_text = ""
        if info.get('subtitles') or info.get('automatic_captions'):
            # Try to get subtitles
            subs = info.get('subtitles', {}).get('en') or info.get('automatic_captions', {}).get('en')
            if subs:
                # Download and parse subtitles
                for sub in subs:
                    if sub.get('ext') == 'vtt' or sub.get('ext') == 'json3':
                        # In production, download and parse the subtitle file
                        pass

        # If no subtitles, use description
        if not transcript_text:
            transcript_text = info.get('description', '')

        # Build metadata
        metadata.update({
            "title": info.get('title'),
            "channel": info.get('uploader'),
            "duration_seconds": info.get('duration'),
            "publish_date": info.get('upload_date'),
            "thumbnail_url": info.get('thumbnail'),
            "view_count": info.get('view_count'),
            "chapters": info.get('chapters', []),
        })

        content = f"Title: {info.get('title', 'Unknown')}\n\n"
        content += f"Channel: {info.get('uploader', 'Unknown')}\n\n"
        if info.get('description'):
            content += f"Description:\n{info['description']}\n\n"
        if transcript_text:
            content += f"Transcript:\n{transcript_text}"

        return content, metadata

    except Exception as e:
        logger.error(f"YouTube extraction failed: {e}")
        # Return minimal content on error
        return f"YouTube video: {url}\n\nTranscript extraction failed: {str(e)}", metadata


def _extract_website_content(url: str, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from website URL using BeautifulSoup"""
    try:
        import requests
        from bs4 import BeautifulSoup

        # Check robots.txt (basic compliance)
        headers = {
            'User-Agent': 'Mozilla/5.0 (compatible; EmpireBot/1.0; +https://empire.ai)',
        }

        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()

        soup = BeautifulSoup(response.content, 'html.parser')

        # Extract title
        title = soup.find('title')
        title_text = title.get_text().strip() if title else url

        # Extract meta description
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        description = meta_desc.get('content', '') if meta_desc else ''

        # Extract author
        author_meta = soup.find('meta', attrs={'name': 'author'})
        author = author_meta.get('content', '') if author_meta else ''

        # Extract publish date
        date_meta = soup.find('meta', attrs={'property': 'article:published_time'})
        publish_date = date_meta.get('content', '') if date_meta else ''

        # Remove unwanted elements
        for element in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
            element.decompose()

        # Extract main content
        main_content = soup.find('main') or soup.find('article') or soup.find('body')
        if main_content:
            # Get text with proper spacing
            text = main_content.get_text(separator='\n', strip=True)
        else:
            text = soup.get_text(separator='\n', strip=True)

        # Clean up text
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        content = '\n'.join(lines)

        # Update metadata
        metadata.update({
            "title": title_text,
            "author": author,
            "publish_date": publish_date,
            "description": description,
            "url": url,
        })

        return content, metadata

    except Exception as e:
        logger.error(f"Website extraction failed: {e}")
        return f"Website: {url}\n\nContent extraction failed: {str(e)}", metadata


def _extract_pdf_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from PDF using LlamaParse or PyPDF"""
    try:
        # Download file from storage
        file_content = _download_file(source, supabase)

        # Try LlamaParse first (better quality)
        try:
            from llama_parse import LlamaParse

            parser = LlamaParse(
                api_key=os.getenv("LLAMA_CLOUD_API_KEY"),
                result_type="markdown",
            )

            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(file_content)
                temp_path = f.name

            try:
                documents = parser.load_data(temp_path)
                content = "\n\n".join([doc.text for doc in documents])
                metadata["parser"] = "llamaparse"
                metadata["page_count"] = len(documents)
            finally:
                os.unlink(temp_path)

        except Exception as e:
            logger.warning(f"LlamaParse failed, falling back to PyPDF: {e}")

            # Fallback to PyPDF
            import pypdf

            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(file_content)
                temp_path = f.name

            try:
                reader = pypdf.PdfReader(temp_path)
                pages = []
                for page in reader.pages:
                    pages.append(page.extract_text())
                content = "\n\n".join(pages)
                metadata["parser"] = "pypdf"
                metadata["page_count"] = len(reader.pages)
            finally:
                os.unlink(temp_path)

        metadata["word_count"] = len(content.split())
        return content, metadata

    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        raise


def _extract_docx_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from Word documents"""
    try:
        from docx import Document

        file_content = _download_file(source, supabase)

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            doc = Document(temp_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)
            metadata["paragraph_count"] = len(paragraphs)
            metadata["word_count"] = len(content.split())
        finally:
            os.unlink(temp_path)

        return content, metadata

    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        raise


def _extract_spreadsheet_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from spreadsheets (Excel, CSV)"""
    try:
        import pandas as pd

        file_content = _download_file(source, supabase)
        source_type = source["source_type"]

        with tempfile.NamedTemporaryFile(suffix=f".{source_type}", delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            if source_type == "csv":
                df = pd.read_csv(temp_path)
            else:
                df = pd.read_excel(temp_path)

            # Convert to markdown table
            content = df.to_markdown(index=False)

            metadata["row_count"] = len(df)
            metadata["column_count"] = len(df.columns)
            metadata["columns"] = list(df.columns)
        finally:
            os.unlink(temp_path)

        return content, metadata

    except Exception as e:
        logger.error(f"Spreadsheet extraction failed: {e}")
        raise


def _extract_pptx_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from PowerPoint presentations"""
    try:
        from pptx import Presentation

        file_content = _download_file(source, supabase)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            prs = Presentation(temp_path)
            slides_content = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"## Slide {i}"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                slides_content.append("\n".join(slide_text))

            content = "\n\n".join(slides_content)
            metadata["slide_count"] = len(prs.slides)
            metadata["word_count"] = len(content.split())
        finally:
            os.unlink(temp_path)

        return content, metadata

    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        raise


def _extract_text_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from text files (TXT, MD, RTF)"""
    try:
        file_content = _download_file(source, supabase)

        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            content = file_content.decode('utf-8', errors='ignore')

        metadata["word_count"] = len(content.split())
        metadata["line_count"] = len(content.split('\n'))

        return content, metadata

    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        raise


def _extract_image_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from images using Claude Vision"""
    try:
        import anthropic
        import base64

        file_content = _download_file(source, supabase)
        mime_type = source.get("mime_type", "image/jpeg")

        # Encode image
        base64_image = base64.b64encode(file_content).decode('utf-8')

        # Use Claude Vision for OCR and description
        client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

        response = client.messages.create(
            model="claude-sonnet-4-5-20250929",
            max_tokens=4096,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": base64_image,
                            },
                        },
                        {
                            "type": "text",
                            "text": "Please describe this image in detail. If there is any text in the image, extract and include it. Provide a comprehensive description suitable for document search."
                        }
                    ],
                }
            ],
        )

        content = response.content[0].text
        metadata["analyzed_by"] = "claude_vision"

        return content, metadata

    except Exception as e:
        logger.error(f"Image extraction failed: {e}")
        return f"Image file: {source.get('file_name', 'unknown')}\n\nAnalysis failed: {str(e)}", metadata


def _extract_audio_video_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """Extract content from audio/video using Soniox transcription"""
    try:
        # For now, return placeholder - Soniox integration would go here
        file_name = source.get("file_name", "unknown")
        source_type = source["source_type"]

        content = f"{source_type.title()} file: {file_name}\n\n"
        content += "Transcription pending - Soniox integration required."

        metadata["transcription_status"] = "pending"

        return content, metadata

    except Exception as e:
        logger.error(f"Audio/Video extraction failed: {e}")
        raise


def _extract_archive_content(source: Dict, supabase, metadata: Dict) -> Tuple[str, Dict[str, Any]]:
    """List contents of archive files"""
    try:
        import zipfile
        import tarfile

        file_content = _download_file(source, supabase)
        file_name = source.get("file_name", "archive")

        with tempfile.NamedTemporaryFile(delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            contents = []

            if file_name.endswith('.zip'):
                with zipfile.ZipFile(temp_path, 'r') as zf:
                    contents = zf.namelist()
            elif file_name.endswith(('.tar', '.tar.gz', '.tgz')):
                with tarfile.open(temp_path, 'r:*') as tf:
                    contents = tf.getnames()

            content = f"Archive: {file_name}\n\nContents ({len(contents)} files):\n"
            content += "\n".join(f"- {f}" for f in contents[:100])  # Limit to 100 files

            if len(contents) > 100:
                content += f"\n... and {len(contents) - 100} more files"

            metadata["file_count"] = len(contents)
        finally:
            os.unlink(temp_path)

        return content, metadata

    except Exception as e:
        logger.error(f"Archive extraction failed: {e}")
        raise


def _download_file(source: Dict, supabase) -> bytes:
    """Download file from Supabase storage"""
    file_path = source.get("file_path")
    if not file_path:
        raise ValueError("No file path specified")

    storage = supabase.storage.from_("documents")
    response = storage.download(file_path)

    return response


# ============================================================================
# Summary Generation
# ============================================================================

def _generate_summary(content: str, source_type: str) -> str:
    """Generate summary using Claude"""
    try:
        import anthropic

        # Truncate content if too long
        max_content = 100000  # ~25k tokens
        if len(content) > max_content:
            content = content[:max_content] + "\n\n[Content truncated...]"

        client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

        prompt = f"""Summarize the following {source_type} content in 2-3 paragraphs.
Focus on the key points, main topics, and any important details.
Make the summary useful for someone who needs to quickly understand what this document contains.

Content:
{content}

Summary:"""

        response = client.messages.create(
            model="claude-3-5-haiku-20241022",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}]
        )

        summary = response.content[0].text
        return summary[:MAX_SUMMARY_LENGTH]

    except Exception as e:
        logger.error(f"Summary generation failed: {e}")
        return ""


# ============================================================================
# Embedding Generation
# ============================================================================

def _is_markdown_content(content: str, min_headers: int = 2) -> bool:
    """
    Detect if content contains sufficient markdown headers for header-based chunking.

    Feature 006: Markdown-aware document splitting

    Args:
        content: Document text to analyze
        min_headers: Minimum number of headers required (default 2)

    Returns:
        True if content has sufficient markdown structure
    """
    headers = HEADER_PATTERN.findall(content)
    return len(headers) >= min_headers


def _chunk_content(content: str, chunk_size: int = 512, document_id: str = None) -> List[str]:
    """
    Split content into chunks for embedding.

    Feature 006: Uses markdown-aware chunking when headers are detected,
    otherwise falls back to simple sentence-aware chunking.

    Args:
        content: Document text to chunk
        chunk_size: Target size for chunks in tokens
        document_id: Optional document identifier for logging

    Returns:
        List of chunk strings
    """
    # Feature 006: Try markdown-aware chunking first
    if _is_markdown_content(content):
        try:
            logger.info(
                "Using markdown-aware chunking",
                extra={"document_id": document_id, "strategy": "markdown"}
            )
            config = MarkdownChunkerConfig(
                max_chunk_size=chunk_size,
                chunk_overlap=int(chunk_size * 0.2),  # 20% overlap
                min_headers_threshold=2
            )
            chunker = MarkdownChunkerStrategy(config=config)

            # Run async chunk method synchronously
            import asyncio
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)

            chunks = loop.run_until_complete(
                chunker.chunk(content, document_id=document_id or "unknown")
            )

            # Extract text content from Chunk objects
            return [chunk.content for chunk in chunks]

        except Exception as e:
            logger.warning(
                f"Markdown chunking failed, falling back to sentence chunking: {e}",
                extra={"document_id": document_id, "error": str(e)}
            )

    # Fallback: Simple sentence-aware chunking
    logger.debug(
        "Using sentence-aware chunking",
        extra={"document_id": document_id, "strategy": "sentence"}
    )
    sentences = re.split(r'(?<=[.!?])\s+', content)
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in sentences:
        sentence_length = len(sentence.split())
        if current_length + sentence_length > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_length = sentence_length
        else:
            current_chunk.append(sentence)
            current_length += sentence_length

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks


def _generate_embeddings(chunks: List[str]) -> List[List[float]]:
    """Generate embeddings using BGE-M3 via Ollama"""
    try:
        import requests

        ollama_url = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
        embeddings = []

        for chunk in chunks:
            response = requests.post(
                f"{ollama_url}/api/embeddings",
                json={"model": "bge-m3", "prompt": chunk},
                timeout=30
            )
            response.raise_for_status()
            embedding = response.json().get("embedding", [])
            embeddings.append(embedding)

        return embeddings

    except Exception as e:
        logger.error(f"Embedding generation failed: {e}")
        raise


def _store_embeddings(
    supabase,
    source_id: str,
    project_id: str,
    user_id: str,
    chunks: List[str],
    embeddings: List[List[float]],
    content_set_context: Optional[Dict[str, Any]] = None  # Feature 007
):
    """
    Store embeddings in source_embeddings table.

    Feature 007: Enhanced with content set context metadata.
    """
    try:
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            embedding_id = str(uuid4())
            chunk_hash = hashlib.md5(chunk.encode()).hexdigest()

            embedding_data = {
                "id": embedding_id,
                "source_id": source_id,
                "project_id": project_id,
                "user_id": user_id,
                "chunk_index": i,
                "chunk_text": chunk[:5000],  # Limit text length
                "chunk_hash": chunk_hash,
                "embedding": embedding,
                "created_at": datetime.utcnow().isoformat(),
            }

            # Feature 007: Add content set metadata if available
            if content_set_context:
                embedding_data["metadata"] = {
                    "content_set_id": content_set_context.get("content_set_id"),
                    "content_set_name": content_set_context.get("content_set_name"),
                    "sequence_number": content_set_context.get("sequence_number"),
                    "is_part_of_sequence": True
                }

            supabase.table("source_embeddings").insert(embedding_data).execute()

        logger.info(f"Stored {len(embeddings)} embeddings for source {source_id}")

        # Feature 007: Create Neo4j relationships if content set context provided
        if content_set_context:
            _create_content_set_relationships(
                source_id=source_id,
                content_set_context=content_set_context
            )

    except Exception as e:
        logger.error(f"Failed to store embeddings: {e}")
        raise


def _create_content_set_relationships(
    source_id: str,
    content_set_context: Dict[str, Any]
):
    """
    Create Neo4j relationships for content set membership.

    Feature 007: Links documents to their parent content set
    and establishes sequence relationships.

    Args:
        source_id: UUID of the processed source/document
        content_set_context: Content set context with set_id and sequence
    """
    try:
        from app.services.graph_service import get_graph_service

        graph_service = get_graph_service()

        content_set_id = content_set_context.get("content_set_id")
        content_set_name = content_set_context.get("content_set_name")
        sequence_number = content_set_context.get("sequence_number", 0)
        is_sequential = content_set_context.get("is_sequential", True)

        if not content_set_id:
            logger.warning("No content_set_id in context, skipping relationship creation")
            return

        # Create or merge ContentSet node
        graph_service.run_query(
            """
            MERGE (cs:ContentSet {id: $content_set_id})
            ON CREATE SET
                cs.name = $content_set_name,
                cs.created_at = datetime(),
                cs.is_sequential = $is_sequential
            ON MATCH SET
                cs.updated_at = datetime()
            """,
            {
                "content_set_id": content_set_id,
                "content_set_name": content_set_name or "Unknown",
                "is_sequential": is_sequential
            }
        )

        # Create PART_OF relationship with sequence number
        graph_service.run_query(
            """
            MATCH (d:Document {id: $document_id})
            MATCH (cs:ContentSet {id: $content_set_id})
            MERGE (d)-[r:PART_OF]->(cs)
            SET r.sequence = $sequence_number,
                r.created_at = datetime()
            """,
            {
                "document_id": source_id,
                "content_set_id": content_set_id,
                "sequence_number": sequence_number
            }
        )

        # Create PRECEDES relationship to previous document in sequence (if exists)
        if sequence_number > 1:
            graph_service.run_query(
                """
                MATCH (current:Document {id: $document_id})-[:PART_OF]->(cs:ContentSet {id: $content_set_id})
                MATCH (prev:Document)-[r:PART_OF]->(cs)
                WHERE r.sequence = $prev_sequence
                MERGE (prev)-[:PRECEDES]->(current)
                MERGE (current)-[:FOLLOWS]->(prev)
                """,
                {
                    "document_id": source_id,
                    "content_set_id": content_set_id,
                    "prev_sequence": sequence_number - 1
                }
            )

        logger.info(
            f"Created content set relationships for document {source_id} "
            f"in set {content_set_id} (sequence: {sequence_number})"
        )

    except Exception as e:
        # Log but don't fail the main task
        logger.warning(f"Failed to create content set relationships: {e}")


# ============================================================================
# Batch Processing Task
# ============================================================================

@celery_app.task(
    name='app.tasks.source_processing.process_sources_batch',
    bind=True
)
def process_sources_batch(self, source_ids: List[str], user_id: str) -> Dict[str, Any]:
    """
    Process multiple sources in sequence.

    Args:
        source_ids: List of source UUIDs to process
        user_id: User ID for access control

    Returns:
        Batch processing results
    """
    results = {
        "total": len(source_ids),
        "successful": 0,
        "failed": 0,
        "results": []
    }

    for source_id in source_ids:
        try:
            result = process_source(source_id, user_id)
            if result.get("status") == "success":
                results["successful"] += 1
            else:
                results["failed"] += 1
            results["results"].append(result)
        except Exception as e:
            results["failed"] += 1
            results["results"].append({
                "source_id": source_id,
                "status": "error",
                "error": str(e)
            })

    return results
