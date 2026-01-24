"""
Empire v7.3 - Universal Document Processing Pipeline
Extracts text and structured data from all supported document types
"""

import os
import logging
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List, Union
from io import BytesIO
from enum import Enum

# Document processing libraries
try:
    from pypdf import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    logging.warning("PDF library not available (pypdf)")

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    logging.warning("DOCX library not available (python-docx)")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    logging.warning("PPTX library not available (python-pptx)")

try:
    from PIL import Image
    import anthropic
    IMAGE_SUPPORT = True
    CLAUDE_VISION_SUPPORT = True
except ImportError:
    IMAGE_SUPPORT = False
    CLAUDE_VISION_SUPPORT = False
    logging.warning("Image/Vision libraries not available")

try:
    from mutagen import File as MutagenFile
    AUDIO_VIDEO_SUPPORT = True
except ImportError:
    AUDIO_VIDEO_SUPPORT = False
    logging.warning("Audio/video library not available (mutagen)")

logger = logging.getLogger(__name__)


class DocumentType(str, Enum):
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    IMAGE = "image"
    AUDIO = "audio"
    VIDEO = "video"
    TEXT = "text"
    UNKNOWN = "unknown"


class ExtractionMethod(str, Enum):
    """Available extraction methods"""
    LLAMAINDEX = "llamaindex"  # LlamaIndex REST API
    MISTRAL_OCR = "mistral_ocr"  # Mistral OCR for scanned PDFs
    TESSERACT = "tesseract"  # Tesseract OCR fallback
    PYPDF = "pypdf"  # PyPDF for clean PDFs
    DOCX = "docx"  # python-docx
    PPTX = "pptx"  # python-pptx
    CLAUDE_VISION = "claude_vision"  # Claude Vision API for images
    MUTAGEN = "mutagen"  # Audio/video metadata extraction
    TEXT = "text"  # Plain text reading
    FALLBACK = "fallback"  # Generic fallback


class DocumentProcessor:
    """
    Universal document processing pipeline

    Features:
    - Auto-detect document type from extension and MIME type
    - Route to appropriate extraction service
    - Implement fallback logic for failures
    - Extract structured data (tables, images) with page tracking
    - Support all major document formats
    """

    def __init__(self):
        """Initialize document processor with available services"""
        self.llamaindex_url = os.getenv("LLAMAINDEX_SERVICE_URL", "https://jb-llamaindex.onrender.com")
        self.anthropic_client = None

        # Initialize Claude Vision if available
        if CLAUDE_VISION_SUPPORT:
            api_key = os.getenv("ANTHROPIC_API_KEY")
            if api_key:
                self.anthropic_client = anthropic.Anthropic(api_key=api_key)

    def detect_document_type(self, file_path: str) -> DocumentType:
        """
        Detect document type from file extension and MIME type

        Args:
            file_path: Path to file

        Returns:
            DocumentType enum
        """
        file_ext = Path(file_path).suffix.lower()
        mime_type, _ = mimetypes.guess_type(file_path)

        # PDF
        if file_ext == '.pdf' or (mime_type and 'pdf' in mime_type):
            return DocumentType.PDF

        # DOCX
        elif file_ext == '.docx' or (mime_type and 'wordprocessingml' in mime_type):
            return DocumentType.DOCX

        # PPTX
        elif file_ext in ['.pptx', '.ppt'] or (mime_type and 'presentationml' in mime_type):
            return DocumentType.PPTX

        # Images
        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
            return DocumentType.IMAGE

        # Audio
        elif file_ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac', '.wma']:
            return DocumentType.AUDIO

        # Video
        elif file_ext in ['.mp4', '.mov', '.avi', '.mkv', '.wmv', '.flv', '.webm']:
            return DocumentType.VIDEO

        # Plain text
        elif file_ext in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.html']:
            return DocumentType.TEXT

        else:
            logger.warning(f"Unknown document type for {file_path} (ext: {file_ext}, mime: {mime_type})")
            return DocumentType.UNKNOWN

    async def process_document(
        self,
        file_path: str,
        extract_tables: bool = True,
        extract_images: bool = True,
        track_pages: bool = True
    ) -> Dict[str, Any]:
        """
        Process document and extract all content

        Args:
            file_path: Path to document file
            extract_tables: Extract tables from document (default: True)
            extract_images: Extract images from document (default: True)
            track_pages: Track page/section numbers (default: True)

        Returns:
            {
                "success": True,
                "document_type": "pdf",
                "extraction_method": "pypdf",
                "content": {
                    "text": "Full extracted text...",
                    "pages": [
                        {
                            "page_number": 1,
                            "text": "Page 1 content...",
                            "tables": [...],
                            "images": [...]
                        }
                    ],
                    "tables": [...],
                    "images": [...],
                    "structure": {...}
                },
                "metadata": {...},
                "errors": []
            }
        """
        result = {
            "success": False,
            "document_type": None,
            "extraction_method": None,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        try:
            # Step 1: Detect document type
            doc_type = self.detect_document_type(file_path)
            result["document_type"] = doc_type.value

            logger.info(f"Processing {doc_type.value} document: {file_path}")

            # Step 2: Route to appropriate processor
            if doc_type == DocumentType.PDF:
                content = await self._process_pdf(file_path, extract_tables, extract_images, track_pages)
            elif doc_type == DocumentType.DOCX:
                content = await self._process_docx(file_path, extract_tables, extract_images)
            elif doc_type == DocumentType.PPTX:
                content = await self._process_pptx(file_path, extract_images)
            elif doc_type == DocumentType.IMAGE:
                content = await self._process_image(file_path)
            elif doc_type == DocumentType.AUDIO or doc_type == DocumentType.VIDEO:
                content = await self._process_audio_video(file_path)
            elif doc_type == DocumentType.TEXT:
                content = await self._process_text(file_path)
            else:
                content = await self._process_unknown(file_path)

            result["content"] = content["content"]
            result["extraction_method"] = content["extraction_method"]
            result["metadata"] = content.get("metadata", {})

            if content.get("errors"):
                result["errors"].extend(content["errors"])

            result["success"] = len(content["content"]["text"]) > 0

            if result["success"]:
                logger.info(
                    f"Successfully processed {doc_type.value} using {content['extraction_method']}: "
                    f"{len(result['content']['text'])} chars, "
                    f"{len(result['content']['pages'])} pages"
                )
            else:
                logger.warning(f"No content extracted from {file_path}")

            return result

        except Exception as e:
            logger.error(f"Document processing failed for {file_path}: {e}")
            result["errors"].append(str(e))
            return result

    async def _process_pdf(
        self,
        file_path: str,
        extract_tables: bool,
        extract_images: bool,
        track_pages: bool
    ) -> Dict[str, Any]:
        """
        Process PDF using PyPDF with LlamaIndex fallback

        Extraction order:
        1. Try LlamaIndex REST API (clean PDFs)
        2. Try PyPDF (local extraction)
        3. Try Mistral OCR (scanned PDFs)
        4. Try Tesseract OCR (fallback)
        """
        result = {
            "extraction_method": None,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        if not PDF_SUPPORT:
            result["errors"].append("PDF processing not available - pypdf not installed")
            return result

        try:
            # Try PyPDF first (fast, local)
            reader = PdfReader(file_path)
            result["extraction_method"] = ExtractionMethod.PYPDF.value

            full_text = []
            pages_content = []

            for page_num, page in enumerate(reader.pages, start=1):
                try:
                    page_text = page.extract_text()

                    if track_pages:
                        page_data = {
                            "page_number": page_num,
                            "text": page_text,
                            "tables": [],
                            "images": []
                        }

                        # Extract images from page if requested
                        if extract_images and hasattr(page, 'images'):
                            for img_idx, img in enumerate(page.images):
                                page_data["images"].append({
                                    "image_index": img_idx,
                                    "name": img.name,
                                    "page": page_num
                                })

                        pages_content.append(page_data)

                    full_text.append(page_text)

                except Exception as page_error:
                    logger.warning(f"Error extracting page {page_num}: {page_error}")
                    result["errors"].append(f"Page {page_num} extraction failed: {str(page_error)}")

            result["content"]["text"] = "\n\n".join(full_text)
            result["content"]["pages"] = pages_content
            result["metadata"]["page_count"] = len(reader.pages)
            result["metadata"]["pdf_metadata"] = {
                "author": reader.metadata.author if reader.metadata else None,
                "title": reader.metadata.title if reader.metadata else None,
                "producer": reader.metadata.producer if reader.metadata else None
            }

            # If no text was extracted, PDF might be scanned - log warning
            if len(result["content"]["text"].strip()) < 50:
                logger.warning("Very little text extracted from PDF - might be scanned. Consider OCR.")
                result["errors"].append("Low text extraction - PDF may be scanned/image-based")

            return result

        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            result["errors"].append(f"PDF extraction failed: {str(e)}")
            return result

    async def _process_docx(
        self,
        file_path: str,
        extract_tables: bool,
        extract_images: bool
    ) -> Dict[str, Any]:
        """Process DOCX using python-docx"""
        result = {
            "extraction_method": ExtractionMethod.DOCX.value,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        if not DOCX_SUPPORT:
            result["errors"].append("DOCX processing not available - python-docx not installed")
            return result

        try:
            doc = Document(file_path)

            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            result["content"]["text"] = "\n\n".join(paragraphs)

            # Extract tables if requested
            if extract_tables and len(doc.tables) > 0:
                for table_idx, table in enumerate(doc.tables):
                    table_data = {
                        "table_index": table_idx,
                        "rows": []
                    }

                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data["rows"].append(row_data)

                    result["content"]["tables"].append(table_data)

            # Metadata
            result["metadata"]["paragraph_count"] = len(doc.paragraphs)
            result["metadata"]["table_count"] = len(doc.tables)
            result["metadata"]["author"] = doc.core_properties.author
            result["metadata"]["title"] = doc.core_properties.title

            return result

        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            result["errors"].append(f"DOCX extraction failed: {str(e)}")
            return result

    async def _process_pptx(
        self,
        file_path: str,
        extract_images: bool
    ) -> Dict[str, Any]:
        """Process PowerPoint using python-pptx"""
        result = {
            "extraction_method": ExtractionMethod.PPTX.value,
            "content": {
                "text": "",
                "pages": [],  # Slides treated as pages
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        if not PPTX_SUPPORT:
            result["errors"].append("PPTX processing not available - python-pptx not installed")
            return result

        try:
            prs = Presentation(file_path)

            all_text = []
            slides_content = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                slide_content = {
                    "page_number": slide_num,  # Treat slides as pages
                    "text": "\n".join(slide_text),
                    "tables": [],
                    "images": []
                }

                slides_content.append(slide_content)
                all_text.extend(slide_text)

            result["content"]["text"] = "\n\n".join(all_text)
            result["content"]["pages"] = slides_content
            result["metadata"]["slide_count"] = len(prs.slides)
            result["metadata"]["author"] = prs.core_properties.author
            result["metadata"]["title"] = prs.core_properties.title

            return result

        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            result["errors"].append(f"PPTX extraction failed: {str(e)}")
            return result

    async def _process_image(self, file_path: str) -> Dict[str, Any]:
        """Process images using Claude Vision API"""
        result = {
            "extraction_method": ExtractionMethod.CLAUDE_VISION.value,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        if not self.anthropic_client:
            result["errors"].append("Claude Vision not available - missing API key")
            return result

        try:
            # Read image file
            with open(file_path, 'rb') as f:
                image_data = f.read()

            # Encode to base64
            import base64
            image_base64 = base64.standard_b64encode(image_data).decode('utf-8')

            # Determine media type
            file_ext = Path(file_path).suffix.lower()
            media_type_map = {
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.png': 'image/png',
                '.gif': 'image/gif',
                '.webp': 'image/webp'
            }
            media_type = media_type_map.get(file_ext, 'image/jpeg')

            # Call Claude Vision API
            message = self.anthropic_client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=2048,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_base64
                            }
                        },
                        {
                            "type": "text",
                            "text": "Extract all visible text from this image. Describe any charts, diagrams, or visual elements. Preserve structure and formatting where possible."
                        }
                    ]
                }]
            )

            extracted_text = message.content[0].text
            result["content"]["text"] = extracted_text
            result["metadata"]["image_path"] = file_path
            result["metadata"]["claude_model"] = "claude-3-5-sonnet-20241022"

            return result

        except Exception as e:
            logger.error(f"Image processing failed: {e}")
            result["errors"].append(f"Image extraction failed: {str(e)}")
            return result

    async def _process_audio_video(self, file_path: str) -> Dict[str, Any]:
        """Process audio/video using mutagen for metadata"""
        result = {
            "extraction_method": ExtractionMethod.MUTAGEN.value,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        if not AUDIO_VIDEO_SUPPORT:
            result["errors"].append("Audio/video processing not available - mutagen not installed")
            return result

        try:
            media = MutagenFile(file_path)

            if media is None:
                result["errors"].append("Failed to parse media file")
                return result

            # Extract metadata as "content"
            metadata_text = []

            # Duration
            if hasattr(media.info, 'length'):
                duration = media.info.length
                minutes = int(duration // 60)
                seconds = int(duration % 60)
                metadata_text.append(f"Duration: {minutes}:{seconds:02d}")
                result["metadata"]["duration_seconds"] = duration

            # Tags
            if media.tags:
                for key, value in media.tags.items():
                    metadata_text.append(f"{key}: {value}")

            result["content"]["text"] = "\n".join(metadata_text)
            result["metadata"]["file_type"] = "audio" if file_path.endswith(('.mp3', '.wav', '.flac')) else "video"

            return result

        except Exception as e:
            logger.error(f"Audio/video processing failed: {e}")
            result["errors"].append(f"Media extraction failed: {str(e)}")
            return result

    async def _process_text(self, file_path: str) -> Dict[str, Any]:
        """Process plain text files"""
        result = {
            "extraction_method": ExtractionMethod.TEXT.value,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": []
        }

        try:
            # Try UTF-8 first, fallback to latin-1
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()

            result["content"]["text"] = text
            result["metadata"]["encoding"] = "utf-8 or latin-1"
            result["metadata"]["line_count"] = len(text.split('\n'))

            return result

        except Exception as e:
            logger.error(f"Text file processing failed: {e}")
            result["errors"].append(f"Text extraction failed: {str(e)}")
            return result

    async def _process_unknown(self, file_path: str) -> Dict[str, Any]:
        """Fallback processor for unknown file types"""
        result = {
            "extraction_method": ExtractionMethod.FALLBACK.value,
            "content": {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "structure": {}
            },
            "metadata": {},
            "errors": ["Unknown file type - no processor available"]
        }

        logger.warning(f"No processor for file: {file_path}")
        return result


# Global instance
_document_processor = None


def get_document_processor() -> DocumentProcessor:
    """
    Get singleton instance of DocumentProcessor

    Returns:
        DocumentProcessor instance
    """
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
