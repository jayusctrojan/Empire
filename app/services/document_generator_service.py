"""
Empire v7.3 - Document Generator Service
Phase 2: Document Generation Pipeline

Generates DOCX, XLSX, PPTX, and PDF files from structured ContentBlocks
produced by the Output Architect. Documents are uploaded to B2 storage.

Pipeline: OutputArchitect → ContentBlocks → DocumentGenerator → B2

Uses asyncio.to_thread() for CPU-bound document generation to avoid
blocking the event loop.
"""

import asyncio
import structlog
import re
from io import BytesIO
from dataclasses import dataclass
from enum import Enum
from typing import List, Optional
from datetime import datetime, timezone

from app.services.output_architect_service import ContentBlock

logger = structlog.get_logger(__name__)


# ============================================================================
# Data Models
# ============================================================================

class DocumentFormat(str, Enum):
    """Supported document formats."""
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    PDF = "pdf"
    MARKDOWN = "md"


MIME_TYPES = {
    DocumentFormat.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    DocumentFormat.XLSX: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    DocumentFormat.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    DocumentFormat.PDF: "application/pdf",
    DocumentFormat.MARKDOWN: "text/markdown",
}


@dataclass
class GeneratedDocument:
    """Result of document generation."""
    format: DocumentFormat
    filename: str
    content_bytes: bytes
    size_bytes: int
    mime_type: str
    storage_url: Optional[str] = None
    storage_path: Optional[str] = None
    preview_markdown: Optional[str] = None


# ============================================================================
# Service
# ============================================================================

class DocumentGeneratorService:
    """
    Generates DOCX, XLSX, PPTX files from structured ContentBlocks.
    All generation runs in asyncio.to_thread() to avoid blocking.
    """

    async def generate(
        self,
        content_blocks: List[ContentBlock],
        format: DocumentFormat,
        title: str,
        summary: Optional[str] = None,
    ) -> GeneratedDocument:
        """
        Generate a document from content blocks.

        Args:
            content_blocks: Structured content from Output Architect
            format: Target document format
            title: Document title
            summary: Optional summary for metadata

        Returns:
            GeneratedDocument with bytes and metadata
        """
        generators = {
            DocumentFormat.DOCX: self._generate_docx,
            DocumentFormat.XLSX: self._generate_xlsx,
            DocumentFormat.PPTX: self._generate_pptx,
            DocumentFormat.PDF: self._generate_pdf,
            DocumentFormat.MARKDOWN: self._generate_markdown,
        }

        generator = generators.get(format)
        if not generator:
            raise ValueError(f"Unsupported format: {format}")

        # Run CPU-bound generation in thread pool
        content_bytes = await asyncio.to_thread(
            generator, content_blocks, title, summary
        )

        # Sanitize filename
        safe_title = re.sub(r'[^\w\s-]', '', title)[:80].strip()
        safe_title = re.sub(r'\s+', '_', safe_title)
        if not safe_title:
            safe_title = "document"
        filename = f"{safe_title}.{format.value}"

        # Generate preview markdown (first ~500 chars of formatted content)
        preview = self._generate_preview(content_blocks)

        return GeneratedDocument(
            format=format,
            filename=filename,
            content_bytes=content_bytes,
            size_bytes=len(content_bytes),
            mime_type=MIME_TYPES[format],
            preview_markdown=preview,
        )

    async def upload_to_storage(
        self,
        document: GeneratedDocument,
        user_id: str,
        session_id: str,
    ) -> GeneratedDocument:
        """
        Upload generated document to B2 storage.

        Args:
            document: Generated document to upload
            user_id: Owner user ID
            session_id: CKO session that produced this document

        Returns:
            Document updated with storage_url and storage_path
        """
        from app.services.b2_storage import get_b2_service, B2Folder

        b2 = get_b2_service()
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        storage_filename = f"{user_id}/{session_id}/{timestamp}_{document.filename}"

        result = await b2.upload_file(
            file_data=BytesIO(document.content_bytes),
            filename=storage_filename,
            folder=B2Folder.ARTIFACTS,
            content_type=document.mime_type,
            metadata={
                "user_id": user_id,
                "session_id": session_id,
                "format": document.format.value,
            },
        )

        document.storage_url = result.get("url", "")
        document.storage_path = result.get("file_name", "")

        logger.info(
            "Document uploaded to B2",
            filename=document.filename,
            size_bytes=document.size_bytes,
            format=document.format.value,
        )

        return document

    # ========================================================================
    # Format-specific generators (run in thread pool)
    # ========================================================================

    def _generate_docx(
        self,
        blocks: List[ContentBlock],
        title: str,
        summary: Optional[str],
    ) -> bytes:
        """Generate a DOCX document from content blocks."""
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Summary / subtitle
        if summary:
            subtitle = doc.add_paragraph(summary)
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
            subtitle.style = doc.styles['Subtitle']

        # Date
        date_para = doc.add_paragraph(
            datetime.now(timezone.utc).strftime("%B %d, %Y")
        )
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph("")  # Spacer

        # Content blocks
        for block in blocks:
            if block.type == "heading":
                level = min(block.metadata.get("level", 2), 3)
                doc.add_heading(block.content, level=level)

            elif block.type == "paragraph":
                doc.add_paragraph(block.content)

            elif block.type == "table":
                headers = block.metadata.get("headers", [])
                rows = block.metadata.get("rows", [])
                if headers:
                    num_cols = len(headers)
                    table = doc.add_table(rows=1 + len(rows), cols=num_cols)
                    table.style = 'Table Grid'

                    # Header row
                    for i, header in enumerate(headers):
                        cell = table.rows[0].cells[i]
                        cell.text = str(header)
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.bold = True

                    # Data rows
                    for row_idx, row_data in enumerate(rows):
                        for col_idx, cell_val in enumerate(row_data):
                            if col_idx < num_cols:
                                table.rows[row_idx + 1].cells[col_idx].text = str(cell_val)

                    doc.add_paragraph("")  # Spacer after table

            elif block.type == "code":
                code_para = doc.add_paragraph()
                code_run = code_para.add_run(block.content.rstrip())
                code_run.font.name = 'Courier New'
                code_run.font.size = Pt(9)

            elif block.type == "list":
                items = block.content.strip().split("\n")
                for item in items:
                    item_text = item.strip()
                    if item_text:
                        doc.add_paragraph(item_text, style='List Bullet')

        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _generate_xlsx(
        self,
        blocks: List[ContentBlock],
        title: str,
        summary: Optional[str],
    ) -> bytes:
        """Generate an XLSX spreadsheet from content blocks."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name limit

        # Styles
        header_font = Font(bold=True, size=12, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        title_font = Font(bold=True, size=16)
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin'),
        )

        current_row = 1

        # Title
        ws.cell(row=current_row, column=1, value=title).font = title_font
        current_row += 1

        if summary:
            ws.cell(row=current_row, column=1, value=summary)
            current_row += 1

        ws.cell(
            row=current_row, column=1,
            value=f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
        )
        current_row += 2

        # Process content blocks — tables get proper spreadsheet treatment
        for block in blocks:
            if block.type == "table":
                headers = block.metadata.get("headers", [])
                rows = block.metadata.get("rows", [])

                if headers:
                    for col_idx, header in enumerate(headers, 1):
                        cell = ws.cell(row=current_row, column=col_idx, value=str(header))
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.alignment = Alignment(horizontal='center')
                        cell.border = thin_border
                    current_row += 1

                    for row_data in rows:
                        for col_idx, val in enumerate(row_data, 1):
                            cell = ws.cell(row=current_row, column=col_idx, value=str(val))
                            cell.border = thin_border
                        current_row += 1

                    current_row += 1  # Gap after table

            elif block.type == "heading":
                level = block.metadata.get("level", 2)
                cell = ws.cell(row=current_row, column=1, value=block.content)
                cell.font = Font(bold=True, size=max(16 - level * 2, 10))
                current_row += 1

            elif block.type == "paragraph":
                ws.cell(row=current_row, column=1, value=block.content)
                current_row += 1

            elif block.type == "list":
                items = block.content.strip().split("\n")
                for item in items:
                    item_text = item.strip()
                    if item_text:
                        ws.cell(row=current_row, column=1, value=f"  - {item_text}")
                        current_row += 1

        # Auto-fit column widths (approximate)
        for col in ws.columns:
            max_length = 0
            col_letter = col[0].column_letter
            for cell in col:
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_length + 4, 60)

        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def _generate_pptx(
        self,
        blocks: List[ContentBlock],
        title: str,
        summary: Optional[str],
    ) -> bytes:
        """Generate a PPTX presentation from content blocks."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()

        # Title slide
        title_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if 1 in slide.placeholders:
            slide.placeholders[1].text = summary or datetime.now(timezone.utc).strftime("%B %d, %Y")

        # Group blocks into slides: each heading starts a new slide
        current_slide = None
        current_body = None

        for block in blocks:
            if block.type == "heading":
                # New slide for each heading
                content_layout = prs.slide_layouts[1]  # Title and Content
                current_slide = prs.slides.add_slide(content_layout)
                current_slide.shapes.title.text = block.content
                if 1 in current_slide.placeholders:
                    current_body = current_slide.placeholders[1].text_frame
                    current_body.clear()
                else:
                    current_body = None

            elif block.type == "table":
                # Create a dedicated slide for the table
                blank_layout = prs.slide_layouts[5]  # Blank
                table_slide = prs.slides.add_slide(blank_layout)

                headers = block.metadata.get("headers", [])
                rows_data = block.metadata.get("rows", [])
                num_rows = len(rows_data) + 1  # +1 for header
                num_cols = len(headers) if headers else 1

                if num_rows > 0 and num_cols > 0:
                    x, y = Inches(0.5), Inches(1.0)
                    cx = Inches(9.0)
                    cy = Inches(0.4 * min(num_rows, 15))

                    shape = table_slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
                    table = shape.table

                    # Headers
                    for col_idx, header in enumerate(headers):
                        cell = table.cell(0, col_idx)
                        cell.text = str(header)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(11)

                    # Data
                    for row_idx, row_data in enumerate(rows_data):
                        for col_idx, val in enumerate(row_data):
                            if col_idx < num_cols:
                                table.cell(row_idx + 1, col_idx).text = str(val)

            elif current_body is not None:
                if block.type == "paragraph":
                    p = current_body.add_paragraph()
                    p.text = block.content
                    p.font.size = Pt(14)

                elif block.type == "list":
                    items = block.content.strip().split("\n")
                    for item in items:
                        item_text = item.strip()
                        if item_text:
                            p = current_body.add_paragraph()
                            p.text = item_text
                            p.level = 0
                            p.font.size = Pt(14)

                elif block.type == "code":
                    p = current_body.add_paragraph()
                    p.text = block.content.rstrip()
                    p.font.name = 'Courier New'
                    p.font.size = Pt(10)

            else:
                # No heading yet — create a generic content slide
                content_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(content_layout)
                current_slide.shapes.title.text = title
                current_body = current_slide.placeholders[1].text_frame
                current_body.clear()

                p = current_body.add_paragraph()
                p.text = block.content if hasattr(block, 'content') else ""
                p.font.size = Pt(14)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _generate_pdf(
        self,
        blocks: List[ContentBlock],
        title: str,
        summary: Optional[str],
    ) -> bytes:
        """Generate a PDF using the existing PDFReportGenerator.

        The PDFReportGenerator.generate() takes (title, content_markdown, metadata)
        so we convert ContentBlocks back to markdown and pass through.
        """
        from app.services.pdf_report_generator import (
            PDFReportGenerator,
            REPORTLAB_AVAILABLE,
        )

        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("reportlab not installed — cannot generate PDF")

        # Convert content blocks to markdown for the PDF generator
        md_bytes = self._generate_markdown(blocks, title, summary)
        markdown_content = md_bytes.decode("utf-8")

        # Strip the title/header since PDFReportGenerator adds its own cover page
        # Remove the "# Title" line and metadata lines at the top
        lines = markdown_content.split("\n")
        content_start = 0
        for i, line in enumerate(lines):
            if line.strip() == "---":
                content_start = i + 1
                break
        content_body = "\n".join(lines[content_start:]).strip()

        generator = PDFReportGenerator()
        metadata = {}
        if summary:
            metadata["subtitle"] = summary
        pdf_bytes = generator.generate(title=title, content=content_body, metadata=metadata)
        return pdf_bytes

    def _generate_markdown(
        self,
        blocks: List[ContentBlock],
        title: str,
        summary: Optional[str],
    ) -> bytes:
        """Generate a Markdown document from content blocks."""
        lines = [f"# {title}\n"]

        if summary:
            lines.append(f"*{summary}*\n")

        lines.append(f"*Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}*\n")
        lines.append("---\n")

        for block in blocks:
            if block.type == "heading":
                level = block.metadata.get("level", 2)
                prefix = "#" * min(level, 4)
                lines.append(f"\n{prefix} {block.content}\n")

            elif block.type == "paragraph":
                lines.append(f"\n{block.content}\n")

            elif block.type == "table":
                headers = block.metadata.get("headers", [])
                rows = block.metadata.get("rows", [])
                if headers:
                    lines.append("\n| " + " | ".join(str(h) for h in headers) + " |")
                    lines.append("| " + " | ".join("---" for _ in headers) + " |")
                    for row in rows:
                        lines.append("| " + " | ".join(str(v) for v in row) + " |")
                    lines.append("")

            elif block.type == "code":
                lang = block.metadata.get("language", "")
                lines.append(f"\n```{lang}")
                lines.append(block.content.rstrip())
                lines.append("```\n")

            elif block.type == "list":
                items = block.content.strip().split("\n")
                lines.append("")
                for item in items:
                    item_text = item.strip()
                    if item_text:
                        lines.append(f"- {item_text}")
                lines.append("")

        return "\n".join(lines).encode("utf-8")

    # ========================================================================
    # Helpers
    # ========================================================================

    @staticmethod
    def _generate_preview(blocks: List[ContentBlock], max_length: int = 500) -> str:
        """Generate a markdown preview from content blocks."""
        lines = []
        total = 0

        for block in blocks:
            if total >= max_length:
                break

            if block.type == "heading":
                level = block.metadata.get("level", 2)
                prefix = "#" * min(level, 4)
                line = f"{prefix} {block.content}"
            elif block.type == "paragraph":
                line = block.content
            elif block.type == "list":
                items = block.content.strip().split("\n")[:3]
                line = "\n".join(f"- {i.strip()}" for i in items if i.strip())
            elif block.type == "table":
                headers = block.metadata.get("headers", [])
                line = f"[Table: {' | '.join(str(h) for h in headers[:5])}]"
            elif block.type == "code":
                lang = block.metadata.get("language", "")
                line = f"```{lang}\n{block.content[:100]}...\n```"
            else:
                line = block.content[:100] if block.content else ""

            lines.append(line)
            total += len(line)

        preview = "\n\n".join(lines)
        if len(preview) > max_length:
            preview = preview[:max_length] + "..."
        return preview


# ============================================================================
# Singleton
# ============================================================================

_service: Optional[DocumentGeneratorService] = None


def get_document_generator_service() -> DocumentGeneratorService:
    global _service
    if _service is None:
        _service = DocumentGeneratorService()
    return _service
